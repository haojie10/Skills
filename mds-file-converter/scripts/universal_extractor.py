import os
import sys
import argparse
from pathlib import Path

# Try importing required libraries
try:
    import fitz  # PyMuPDF
    import pdfplumber
except ImportError:
    fitz = None
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    from openpyxl_image_loader import SheetImageLoader
except ImportError:
    openpyxl = None
    SheetImageLoader = None

def setup_temp_dirs(base_dir):
    temp_dir = Path(base_dir) / "temp"
    images_dir = temp_dir / "images"
    os.makedirs(images_dir, exist_ok=True)
    return temp_dir, images_dir

def process_pdf(file_path, raw_text_file, images_dir):
    if not fitz or not pdfplumber:
        raise ImportError("Missing dependencies for PDF processing. Run: pip install PyMuPDF pdfplumber")
    
    print(f"[*] Processing PDF: {file_path}")
    doc = fitz.open(file_path)
    
    with open(raw_text_file, "w", encoding="utf-8") as f:
        f.write(f"--- SOURCE: {file_path} ---\n\n")
        
        with pdfplumber.open(file_path) as pdf:
            for page_num in range(len(doc)):
                f.write(f"=== PAGE {page_num + 1} ===\n")
                
                # 1. Text Extraction
                page_text = ""
                if page_num < len(pdf.pages):
                    plumber_page = pdf.pages[page_num]
                    page_text = plumber_page.extract_text() or ""
                
                f.write("[TEXT CONTENT]\n")
                f.write(page_text.strip() + "\n\n")
                
                # 2. Image Extraction
                fitz_page = doc[page_num]
                image_list = fitz_page.get_images(full=True)
                
                # If page has very little text, render the whole page as fallback
                if len(page_text.strip()) < 50:
                    f.write("[PAGE RENDER (Low text detected)]\n")
                    pix = fitz_page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_name = f"page_{page_num + 1}_render.png"
                    img_path = images_dir / img_name
                    pix.save(str(img_path))
                    f.write(f"Rendered Page Image: {img_path.absolute()}\n\n")
                
                f.write("[EXTRACTED IMAGES]\n")
                if image_list:
                    for img_index, img in enumerate(image_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        img_name = f"page_{page_num + 1}_img_{img_index + 1}.{image_ext}"
                        img_path = images_dir / img_name
                        
                        with open(img_path, "wb") as img_file:
                            img_file.write(image_bytes)
                        f.write(f"Image {img_index + 1}: {img_path.absolute()}\n")
                else:
                    f.write("No extractable images found on this page.\n")
                
                f.write("\n" + "="*40 + "\n\n")
    doc.close()


def process_excel(file_path, raw_text_file, images_dir):
    if not openpyxl:
        raise ImportError("Missing dependencies for Excel processing. Run: pip install openpyxl openpyxl-image-loader")
    
    print(f"[*] Processing Excel: {file_path}")
    wb = openpyxl.load_workbook(file_path, data_only=True)
    
    with open(raw_text_file, "w", encoding="utf-8") as f:
        f.write(f"--- SOURCE: {file_path} ---\n\n")
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            f.write(f"=== SHEET: {sheet_name} ===\n")
            
            # Text Extraction
            f.write("[TEXT CONTENT (Row by Row)]\n")
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                # Filter out completely empty rows
                if not any(row): continue
                row_str = " | ".join(str(cell).strip() if cell is not None else "" for cell in row)
                f.write(f"Row {row_idx + 1}: {row_str}\n")
            
            print(f"[*] Extracting images from sheet: {sheet_name}")
            try:
                # Need SheetImageLoader if using openpyxl-image-loader
                image_loader = SheetImageLoader(sheet)
                
                # Iterate through cells to find images
                img_count = 0
                for row in sheet.iter_rows():
                    for cell in row:
                        try:
                            if image_loader.image_in(cell.coordinate):
                                img = image_loader.get(cell.coordinate)
                                img_name = f"sheet_{sheet_name}_cell_{cell.coordinate}.png"
                                img_path = images_dir / img_name
                                img.save(str(img_path))
                                f.write(f"Image at {cell.coordinate}: {img_path.absolute()}\n")
                                img_count += 1
                        except Exception as cell_e:
                            # Skip problematic cells but keep going
                            continue
                if img_count == 0:
                    f.write("No images found in this sheet.\n")
                    
            except Exception as e:
                f.write(f"Warning: Could not extract images properly. Error: {e}\n")
                print(f"[WARNING] Image extraction failed for sheet {sheet_name}: {e}")
            
            f.write("\n" + "="*40 + "\n\n")

def process_ppt(file_path, raw_text_file, images_dir):
    if not Presentation:
        raise ImportError("Missing dependencies for PPT processing. Run: pip install python-pptx")
    
    print(f"[*] Processing PPTX: {file_path}")
    prs = Presentation(file_path)
    
    with open(raw_text_file, "w", encoding="utf-8") as f:
        f.write(f"--- SOURCE: {file_path} ---\n\n")
        
        for slide_idx, slide in enumerate(prs.slides):
            f.write(f"=== SLIDE {slide_idx + 1} ===\n")
            
            f.write("[TEXT CONTENT]\n")
            has_text = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    has_text = True
                    text = shape.text.strip()
                    if text:
                        f.write(text + "\n")
            if not has_text:
                f.write("(No text on this slide)\n")
                
            f.write("\n[EXTRACTED IMAGES]\n")
            img_count = 0
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13: # Picture
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.ext
                    img_name = f"slide_{slide_idx + 1}_img_{shape_idx + 1}.{ext}"
                    img_path = images_dir / img_name
                    
                    with open(img_path, "wb") as img_file:
                        img_file.write(img_bytes)
                    f.write(f"Image {shape_idx + 1}: {img_path.absolute()}\n")
                    img_count += 1
            if img_count == 0:
                f.write("No images found on this slide.\n")
                
            f.write("\n" + "="*40 + "\n\n")


def main():
    parser = argparse.ArgumentParser(description="MDS Universal Resource Extractor: Dumps text and images from PDF/XLSX/PPTX.")
    parser.add_argument("input_file", help="Path to the source document.")
    parser.add_argument("--output_dir", default=".", help="Base directory to create the temp/ folder. Default is current dir.")
    
    args = parser.parse_args()
    input_path = Path(args.input_file)
    
    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
        
    temp_dir, images_dir = setup_temp_dirs(args.output_dir)
    raw_text_file = temp_dir / "raw_data.txt"
    
    ext = input_path.suffix.lower()
    
    try:
        if ext == '.pdf':
            process_pdf(input_path, raw_text_file, images_dir)
        elif ext in ['.xlsx', '.xlsm']:
            process_excel(input_path, raw_text_file, images_dir)
        elif ext in ['.pptx', '.ppt']:
            process_ppt(input_path, raw_text_file, images_dir)
        else:
            print(f"Error: Unsupported file format {ext}. Supported: .pdf, .xlsx, .pptx")
            sys.exit(1)
            
        print(f"\n[SUCCESS] Extraction complete.")
        print(f"Text data written to: {raw_text_file.absolute()}")
        print(f"Images extracted to: {images_dir.absolute()}")
        
    except Exception as e:
        print(f"\n[ERROR] Extraction failed: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
